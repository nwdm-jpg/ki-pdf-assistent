"""Formatunabhängige Dokumentverarbeitung: Erkennt den Dateityp, extrahiert
Text und teilt ihn in Chunks für die semantische Suche.

Ersetzt das frühere `pdf_verarbeitung.py` (PDF-only). Jeder Dateityp hat
einen eigenen kleinen Parser, der die Datei in "Einheiten" zerlegt -
Dicts {"nummer": int, "anzeige": str, "text": str}, jeweils eine
logische Seite/Folie/Tabellenblatt/Abschnitt. Alle Parser münden in
`dokument_verarbeiten`, das die Einheiten anschließend mit der
gemeinsamen `text_in_chunks_aufteilen`-Logik (unverändert aus dem
früheren `pdf_verarbeitung.py`) in überlappende Text-Chunks teilt.

Neue Formate lassen sich ergänzen, indem eine neue `_..._einheiten`-
Funktion geschrieben und in `PARSER_JE_ENDUNG` eingetragen wird - der
Rest der Pipeline (Chunking, Speicherung, Retrieval, Quellenformatierung
über `quellen.py`) bleibt unverändert.
"""

import csv
import io
import re
from pathlib import Path

import openpyxl
from docx import Document
from pptx import Presentation
from pypdf import PdfReader


CHUNK_GROESSE = 1000
CHUNK_UEBERLAPPUNG = 150

# Zielgröße für Text-Abschnitte bei paragraphenbasierten Formaten (DOCX,
# TXT, MD) - grob vergleichbar mit einer PDF-Seite, damit die
# nachgelagerte Chunk-Aufteilung auf ähnlich großen Einheiten arbeitet.
ABSCHNITT_ZIELGROESSE = 1500

CSV_ZEILEN_PRO_ABSCHNITT = 50
XLSX_MAX_ZEILEN_PRO_BLATT = 500


class NichtUnterstuetzterDateityp(Exception):
    """Wird geworfen, wenn für eine Dateiendung kein Parser existiert."""


def text_in_chunks_aufteilen(text, chunk_groesse=CHUNK_GROESSE, ueberlappung=CHUNK_UEBERLAPPUNG):
    """Teilt einen Text in überlappende Chunks, bevorzugt an Satzgrenzen.

    Kurze Texte (kürzer als `chunk_groesse`) werden unverändert als
    einzelner Chunk zurückgegeben.
    """
    text = text.strip()

    if not text:
        return []

    if len(text) <= chunk_groesse:
        return [text]

    saetze = re.split(r"(?<=[.!?])\s+", text)

    chunks = []
    aktueller_chunk = ""

    for satz in saetze:
        if aktueller_chunk and len(aktueller_chunk) + len(satz) + 1 > chunk_groesse:
            chunks.append(aktueller_chunk.strip())
            ueberlapp_text = aktueller_chunk[-ueberlappung:]
            aktueller_chunk = f"{ueberlapp_text} {satz}".strip()
        else:
            aktueller_chunk = f"{aktueller_chunk} {satz}".strip()

    if aktueller_chunk.strip():
        chunks.append(aktueller_chunk.strip())

    return chunks


def _bytes_zu_text(datei_bytes):
    """Dekodiert Rohtext mit UTF-8, fällt auf Latin-1 zurück (verlustfrei)."""
    try:
        return datei_bytes.decode("utf-8")
    except UnicodeDecodeError:
        return datei_bytes.decode("latin-1")


def _absaetze_zu_abschnitten(absaetze, zielgroesse=ABSCHNITT_ZIELGROESSE):
    """Gruppiert eine Liste von Absätzen zu Abschnitten mit Zielgröße.

    Genutzt von DOCX (Absätze) sowie TXT und MD (durch Leerzeilen
    getrennte Absätze). Ein neuer Abschnitt beginnt, sobald der aktuelle
    die Zielgröße überschreitet - eine einfache, aber für Retrieval-
    Zwecke ausreichende Heuristik (vgl. PDF-Seiten, die ebenfalls
    unabhängig von ihrem tatsächlichen Inhalt die Chunking-Einheit
    bilden).
    """
    einheiten = []
    aktuelle_teile = []
    aktuelle_groesse = 0
    nummer = 1

    for absatz in absaetze:
        absatz = absatz.strip()

        if not absatz:
            continue

        if aktuelle_teile and aktuelle_groesse + len(absatz) > zielgroesse:
            einheiten.append(
                {"nummer": nummer, "anzeige": str(nummer), "text": "\n".join(aktuelle_teile)}
            )
            nummer += 1
            aktuelle_teile = []
            aktuelle_groesse = 0

        aktuelle_teile.append(absatz)
        aktuelle_groesse += len(absatz)

    if aktuelle_teile:
        einheiten.append(
            {"nummer": nummer, "anzeige": str(nummer), "text": "\n".join(aktuelle_teile)}
        )

    return einheiten


def _pdf_einheiten(datei_bytes):
    reader = PdfReader(io.BytesIO(datei_bytes))
    einheiten = []

    for nummer, seite in enumerate(reader.pages, start=1):
        text = seite.extract_text()

        if text:
            einheiten.append({"nummer": nummer, "anzeige": str(nummer), "text": text})

    return einheiten, "seite"


def _docx_einheiten(datei_bytes):
    dokument = Document(io.BytesIO(datei_bytes))
    absaetze = [absatz.text for absatz in dokument.paragraphs]

    return _absaetze_zu_abschnitten(absaetze), "abschnitt"


def _txt_einheiten(datei_bytes):
    text = _bytes_zu_text(datei_bytes)
    absaetze = re.split(r"\n\s*\n", text)

    return _absaetze_zu_abschnitten(absaetze), "abschnitt"


def _md_einheiten(datei_bytes):
    # Markdown wird wie Text behandelt - Absätze (durch Leerzeilen
    # getrennt) werden zu Abschnitten gruppiert, ohne die
    # Markdown-Syntax gesondert zu interpretieren.
    einheiten, _ = _txt_einheiten(datei_bytes)
    return einheiten, "abschnitt"


def _csv_einheiten(datei_bytes):
    text = _bytes_zu_text(datei_bytes)
    zeilen = list(csv.reader(io.StringIO(text)))

    einheiten = []

    for start in range(0, len(zeilen), CSV_ZEILEN_PRO_ABSCHNITT):
        block = zeilen[start : start + CSV_ZEILEN_PRO_ABSCHNITT]
        text_block = "\n".join(" | ".join(zeile) for zeile in block if any(zeile))

        if text_block.strip():
            nummer = start // CSV_ZEILEN_PRO_ABSCHNITT + 1
            einheiten.append({"nummer": nummer, "anzeige": str(nummer), "text": text_block})

    return einheiten, "abschnitt"


def _xlsx_einheiten(datei_bytes):
    arbeitsmappe = openpyxl.load_workbook(
        io.BytesIO(datei_bytes), data_only=True, read_only=True
    )

    einheiten = []

    for nummer, blatt in enumerate(arbeitsmappe.worksheets, start=1):
        zeilen_text = []

        for zeilenindex, zeile in enumerate(blatt.iter_rows(values_only=True)):
            if zeilenindex >= XLSX_MAX_ZEILEN_PRO_BLATT:
                break

            werte = [str(wert) for wert in zeile if wert is not None]

            if werte:
                zeilen_text.append(" | ".join(werte))

        if zeilen_text:
            einheiten.append(
                {"nummer": nummer, "anzeige": blatt.title, "text": "\n".join(zeilen_text)}
            )

    return einheiten, "tabellenblatt"


def _pptx_einheiten(datei_bytes):
    praesentation = Presentation(io.BytesIO(datei_bytes))
    einheiten = []

    for nummer, folie in enumerate(praesentation.slides, start=1):
        texte = [
            form.text_frame.text
            for form in folie.shapes
            if form.has_text_frame and form.text_frame.text.strip()
        ]

        if texte:
            einheiten.append(
                {"nummer": nummer, "anzeige": str(nummer), "text": "\n".join(texte)}
            )

    return einheiten, "folie"


# Zentrale Zuordnung Dateiendung -> Parserfunktion. Jeder Eintrag hier
# entspricht einem tatsächlich implementierten und getesteten Parser -
# es wird kein Format als unterstützt geführt, für das es keinen
# funktionierenden Parser gibt.
PARSER_JE_ENDUNG = {
    "pdf": _pdf_einheiten,
    "docx": _docx_einheiten,
    "txt": _txt_einheiten,
    "md": _md_einheiten,
    "csv": _csv_einheiten,
    "xlsx": _xlsx_einheiten,
    "pptx": _pptx_einheiten,
}

SUPPORTED_EXTENSIONS = sorted(PARSER_JE_ENDUNG.keys())


def dokument_verarbeiten(dateiname, datei_bytes):
    """Erkennt den Dateityp, extrahiert Text und teilt ihn in Chunks.

    Gibt (chunks, einheiten_anzahl, einheit_typ, dateityp) zurück. Jeder
    Chunk ist ein Dict {"dateiname", "seitennummer", "einheit_typ",
    "einheit_anzeige", "text"} - direkt kompatibel mit `speicher.py`
    (Speicherung) und `quellen.py` (Quellenformatierung). Wirft
    `NichtUnterstuetzterDateityp` bei unbekannter Endung.
    """
    dateityp = Path(dateiname).suffix.lower().lstrip(".")
    parser = PARSER_JE_ENDUNG.get(dateityp)

    if parser is None:
        raise NichtUnterstuetzterDateityp(
            f"Der Dateityp „.{dateityp}“ wird nicht unterstützt."
        )

    einheiten, einheit_typ = parser(datei_bytes)

    chunks = [
        {
            "dateiname": dateiname,
            "seitennummer": einheit["nummer"],
            "einheit_typ": einheit_typ,
            "einheit_anzeige": einheit["anzeige"],
            "text": chunk_text,
        }
        for einheit in einheiten
        for chunk_text in text_in_chunks_aufteilen(einheit["text"])
    ]

    return chunks, len(einheiten), einheit_typ, dateityp
